from pptx import Presentation
from pptx.util import Inches, Pt

# Load the existing presentation
prs = Presentation()

# Titles and contents for the slides
slides_content = [
    ("GTR(유도 조직 재생술)과 뼈 이식재를 사용한 골내낭 결손의 치료: 임상적 결과 및 예후 인자 평가", 
     "Jad Majzoub, Shayan Barootchi, Lorenzo Tavelli, Chin-Wei Wang, Hsun-Liang Chan, Hom-Lay Wang\nJournal of Periodontology, 2020"),
    
    ("서론", 
     "치주염은 복잡한 다인성 질병으로 골내낭 결손을 유발함.\n다양한 치료 옵션이 있지만, GTR은 장기적으로 임상적, 조직학적 재생을 입증함.\n환자 요인, 결손 형태, 수술 기법 등이 GTR의 예후에 영향을 미침."),
    
    ("재료 및 방법", 
     "연구 설계: 회고적 코호트 연구\n포함 기준: 최소 PD 6mm 이상의 골내낭 결손, SRP 및 TBI 후 GTR 치료, 최소 1년 이상의 추적 관찰\n제외 기준: 12개월 미만의 추적 관찰, 치근이개부에 GTR을 받은 경우, 비흡수성 차폐막 사용 등"),
    
    ("결과", 
     "175개의 결손 부위를 평균 5.75년 동안 추적 관찰\n5년 생존율 85%, 10년 생존율 72.7%\n다양한 변수들이 임상 부착 수준(CAL)과 치아 생존에 미치는 영향 평가"),
    
    ("토론", 
     "GTR의 임상적 결과는 매우 긍정적임.\n흡연, 멤브레인 노출 등은 CAL 감소와 유의하게 관련 있음.\n나이, 연간 관리 횟수는 치아 생존에 중요한 예측 변수임."),
    
    ("결론", 
     "GTR은 골내낭 결손 치료에 효과적인 방법임.\n치료 결과는 환자의 유지 관리, 흡연 여부, 초기 결손 특성 등에 크게 영향을 받음.\n지속적인 관리와 흡연 중단이 중요함."),
    
    ("질문 및 답변", 
     "질문을 받고 답변하는 시간")
]

# Create slides with the given content
for title, content in slides_content:
    slide_layout = prs.slide_layouts[1]  # Using the 'Title and Content' layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = content

    # Formatting the content text
    for paragraph in content_placeholder.text_frame.paragraphs:
        paragraph.font.size = Pt(18)

# Add table and figure images to appropriate slides
table_slide_index = 3  # Results slide index
figure_slide_index = 5  # Discussion slide index

# Add table image to the Results slide
table_slide = prs.slides[table_slide_index]
left = Inches(0.5)
top = Inches(3.5)
height = Inches(2)
table_slide.shapes.add_picture('/mnt/data/table_image.png', left, top, height=height)

# Add figure image to the Discussion slide
figure_slide = prs.slides[figure_slide_index]
left = Inches(0.5)
top = Inches(3.5)
height = Inches(2)
figure_slide.shapes.add_picture('/mnt/data/figure_image.png', left, top, height=height)

# Save the modified presentation
output_path = "/mnt/data/Updated_Presentation_with_Images.pptx"
prs.save(output_path)

output_path
